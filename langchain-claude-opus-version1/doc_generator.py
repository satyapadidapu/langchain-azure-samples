"""Document generation module — creates PDF, DOCX, XLSX, PPTX, CSV, JSON, and chart images."""

import json
import csv
import os
import tempfile
from pathlib import Path
from datetime import datetime

from fpdf import FPDF
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

OUTPUT_DIR = Path(tempfile.gettempdir()) / "langchain_generated_docs"
OUTPUT_DIR.mkdir(exist_ok=True)


def _ts():
    return datetime.now().strftime("%Y%m%d_%H%M%S")


# ─── PDF ──────────────────────────────────────────────────────────────────

def _sanitize(text):
    """Replace Unicode chars unsupported by Helvetica with ASCII equivalents."""
    replacements = {
        "\u2014": "--", "\u2013": "-", "\u2018": "'", "\u2019": "'",
        "\u201c": '"', "\u201d": '"', "\u2026": "...", "\u2022": "*",
        "\u00a0": " ", "\u200b": "", "\u2010": "-", "\u2011": "-",
        "\u2012": "-", "\u00b7": ".", "\u2032": "'", "\u2033": '"',
    }
    for uc, asc in replacements.items():
        text = text.replace(uc, asc)
    # Fallback: replace any remaining non-latin1 chars
    return text.encode("latin-1", errors="replace").decode("latin-1")


class _StyledPDF(FPDF):
    """Custom FPDF subclass with header/footer, page numbers, and accent colors."""

    ACCENT = (41, 98, 166)       # Dark blue
    ACCENT_LIGHT = (220, 230, 241)  # Light blue for table stripes
    GRAY = (100, 100, 100)

    def __init__(self, doc_title=""):
        super().__init__()
        self._doc_title = doc_title

    def header(self):
        if self.page_no() == 1:
            return  # Title page has no header
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(*self.GRAY)
        self.cell(0, 8, _sanitize(self._doc_title), align="L")
        self.ln(2)
        # Accent line below header
        self.set_draw_color(*self.ACCENT)
        self.set_line_width(0.4)
        self.line(self.l_margin, self.get_y(), self.w - self.r_margin, self.get_y())
        self.ln(6)

    def footer(self):
        self.set_y(-15)
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(*self.GRAY)
        self.cell(0, 10, f"Page {self.page_no()}/{{nb}}", align="C")


def generate_pdf(content, filename=None):
    filepath = OUTPUT_DIR / (filename or f"document_{_ts()}.pdf")
    title_text = content.get("title", "Document")
    pdf = _StyledPDF(doc_title=title_text)
    pdf.alias_nb_pages()
    pdf.set_auto_page_break(auto=True, margin=20)
    pdf.set_left_margin(20)
    pdf.set_right_margin(20)
    usable_w = pdf.w - 40  # left + right margins

    # ── Title Page ──
    pdf.add_page()
    pdf.ln(50)
    # Accent bar
    pdf.set_fill_color(*_StyledPDF.ACCENT)
    pdf.rect(20, pdf.get_y(), usable_w, 3, "F")
    pdf.ln(12)
    # Title
    pdf.set_font("Helvetica", "B", 28)
    pdf.set_text_color(*_StyledPDF.ACCENT)
    pdf.multi_cell(0, 14, _sanitize(title_text), align="C")
    pdf.ln(6)
    # Subtitle / date
    pdf.set_font("Helvetica", "", 12)
    pdf.set_text_color(*_StyledPDF.GRAY)
    pdf.cell(0, 8, f"Generated on {datetime.now().strftime('%B %d, %Y')}", align="C",
             new_x="LMARGIN", new_y="NEXT")
    pdf.ln(8)
    # Bottom accent bar
    pdf.set_fill_color(*_StyledPDF.ACCENT)
    pdf.rect(20, pdf.get_y(), usable_w, 3, "F")

    # ── Content Pages ──
    for idx, section in enumerate(content.get("sections", [])):
        # Start a new page for each section
        pdf.add_page()

        # Section heading with accent color
        pdf.set_font("Helvetica", "B", 16)
        pdf.set_text_color(*_StyledPDF.ACCENT)
        heading = _sanitize(section.get("heading", ""))
        pdf.cell(0, 10, heading, new_x="LMARGIN", new_y="NEXT")
        # Thin accent line under heading
        pdf.set_draw_color(*_StyledPDF.ACCENT)
        pdf.set_line_width(0.3)
        pdf.line(pdf.l_margin, pdf.get_y(), pdf.l_margin + pdf.get_string_width(heading) + 5, pdf.get_y())
        pdf.ln(6)

        # Body text
        body = section.get("body", "")
        if body:
            pdf.set_font("Helvetica", "", 11)
            pdf.set_text_color(30, 30, 30)
            pdf.multi_cell(0, 6.5, _sanitize(body))
            pdf.ln(4)

        # Bullets (if present)
        for bullet in section.get("bullets", []):
            pdf.set_font("Helvetica", "", 11)
            pdf.set_text_color(30, 30, 30)
            x_start = pdf.get_x()
            pdf.cell(6, 6.5, chr(149))  # bullet char
            pdf.multi_cell(usable_w - 6, 6.5, _sanitize(str(bullet)))

        # Table (if present)
        if "table" in section:
            table = section["table"]
            headers = table.get("headers", [])
            rows = table.get("rows", [])
            if headers:
                pdf.ln(3)
                col_w = usable_w / len(headers)
                # Header row — dark accent fill
                pdf.set_fill_color(*_StyledPDF.ACCENT)
                pdf.set_text_color(255, 255, 255)
                pdf.set_font("Helvetica", "B", 10)
                for h in headers:
                    pdf.cell(col_w, 9, _sanitize(str(h)), border=1, align="C", fill=True)
                pdf.ln()
                # Data rows — alternating stripes
                pdf.set_text_color(30, 30, 30)
                pdf.set_font("Helvetica", "", 10)
                for r_idx, row in enumerate(rows):
                    if r_idx % 2 == 0:
                        pdf.set_fill_color(*_StyledPDF.ACCENT_LIGHT)
                    else:
                        pdf.set_fill_color(255, 255, 255)
                    for cell in row:
                        pdf.cell(col_w, 8, _sanitize(str(cell)), border=1, fill=True)
                    pdf.ln()
                pdf.ln(4)

    pdf.output(str(filepath))
    return str(filepath)


# ─── DOCX ─────────────────────────────────────────────────────────────────

def generate_docx(content, filename=None):
    filepath = OUTPUT_DIR / (filename or f"document_{_ts()}.docx")
    doc = Document()

    title_heading = doc.add_heading(content.get("title", "Document"), level=0)
    title_heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

    for section in content.get("sections", []):
        level = min(section.get("level", 1), 4)
        doc.add_heading(section.get("heading", ""), level=level)

        for item in section.get("content", []):
            t = item.get("type", "paragraph")

            if t == "paragraph":
                doc.add_paragraph(item.get("text", ""))

            elif t == "bullets":
                for bullet in item.get("items", []):
                    doc.add_paragraph(bullet, style="List Bullet")

            elif t == "numbered":
                for numbered in item.get("items", []):
                    doc.add_paragraph(numbered, style="List Number")

            elif t == "table":
                headers = item.get("headers", [])
                rows = item.get("rows", [])
                if headers:
                    tbl = doc.add_table(rows=1 + len(rows), cols=len(headers))
                    tbl.style = "Light Grid Accent 1"
                    for i, h in enumerate(headers):
                        tbl.rows[0].cells[i].text = str(h)
                    for r_idx, row in enumerate(rows):
                        for c_idx, cell in enumerate(row):
                            if c_idx < len(headers):
                                tbl.rows[r_idx + 1].cells[c_idx].text = str(cell)

    doc.save(str(filepath))
    return str(filepath)


# ─── XLSX ─────────────────────────────────────────────────────────────────

def generate_xlsx(content, filename=None):
    filepath = OUTPUT_DIR / (filename or f"document_{_ts()}.xlsx")
    wb = Workbook()
    wb.remove(wb.active)

    hdr_font = Font(bold=True, size=11, color="FFFFFF")
    hdr_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    thin = Border(
        left=Side(style="thin"), right=Side(style="thin"),
        top=Side(style="thin"), bottom=Side(style="thin"),
    )

    sheets = content.get("sheets", [])
    if not sheets:
        sheets = [{"name": "Sheet1", "headers": content.get("headers", []), "rows": content.get("rows", [])}]

    for sd in sheets:
        ws = wb.create_sheet(title=str(sd.get("name", "Sheet1"))[:31])
        headers = sd.get("headers", [])
        rows = sd.get("rows", [])

        for col, h in enumerate(headers, 1):
            c = ws.cell(row=1, column=col, value=str(h))
            c.font = hdr_font
            c.fill = hdr_fill
            c.alignment = Alignment(horizontal="center")
            c.border = thin

        for r_idx, row in enumerate(rows, 2):
            for c_idx, val in enumerate(row, 1):
                c = ws.cell(row=r_idx, column=c_idx)
                try:
                    c.value = float(val) if "." in str(val) else int(val)
                except (ValueError, TypeError):
                    c.value = str(val)
                c.border = thin

        for col in ws.columns:
            mx = max((len(str(cell.value or "")) for cell in col), default=8)
            ws.column_dimensions[col[0].column_letter].width = min(mx + 4, 50)

    wb.save(str(filepath))
    return str(filepath)


# ─── PPTX ─────────────────────────────────────────────────────────────────

def _pptx_set_bg(slide, r, g, b):
    """Set a solid background colour on a slide."""
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def _add_pptx_shape_rect(slide, left, top, width, height, r, g, b):
    """Add a solid-colour rectangle as a decorative shape."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    shape.line.fill.background()  # no border
    return shape


def generate_pptx(content, filename=None, design=None):
    filepath = OUTPUT_DIR / (filename or f"presentation_{_ts()}.pptx")
    prs = Presentation()

    # ── Apply design: use extracted values or defaults ──
    d = design or {}
    colors = d.get("colors", {})
    fonts = d.get("fonts", {})

    sw = d.get("slide_width", 13.33)
    sh = d.get("slide_height", 7.5)
    prs.slide_width = PptxInches(sw)
    prs.slide_height = PptxInches(sh)

    dark_bg = colors.get("dark_bg", (30, 58, 95))
    accent = colors.get("accent", (41, 98, 166))
    accent2 = colors.get("accent2", (86, 156, 214))
    light_bg = colors.get("light_bg", (240, 244, 248))
    title_text_c = colors.get("title_text", (255, 255, 255))
    body_text_c = colors.get("body_text", (40, 40, 40))

    DARK = RGBColor(*dark_bg)
    ACCENT = RGBColor(*accent)
    ACCENT2 = RGBColor(*accent2)
    WHITE = RGBColor(*title_text_c)
    LIGHT_BG = RGBColor(*light_bg)
    DARK_TEXT = RGBColor(*body_text_c)

    title_font_sz = fonts.get("title_size", 40)
    subtitle_font_sz = max(title_font_sz - 20, 16)
    body_font_sz = fonts.get("body_size", 18)
    bullet_font_sz = fonts.get("bullet_size", 12)
    bar_w = d.get("accent_bar_width", 0.4) if d.get("has_accent_bar", True) else 0.4
    bullet_char = d.get("bullet_char", "\u25A0")

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # ── Title slide (custom dark background) ──
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _pptx_set_bg(title_slide, *dark_bg)

    # Accent bar left
    _add_pptx_shape_rect(title_slide, PptxInches(0), PptxInches(0),
                         PptxInches(bar_w), slide_h, *accent)

    # Title text box
    from pptx.util import Emu
    txBox = title_slide.shapes.add_textbox(
        PptxInches(1.5), PptxInches(2.0), PptxInches(10), PptxInches(2.0)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content.get("title", "Presentation")
    p.font.size = PptxPt(title_font_sz)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    subtitle = content.get("subtitle", "")
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = PptxPt(subtitle_font_sz)
        p2.font.color.rgb = ACCENT2
        p2.space_before = PptxPt(14)

    # Bottom accent line
    _add_pptx_shape_rect(title_slide, PptxInches(1.5), PptxInches(4.5),
                         PptxInches(4), PptxInches(0.06), *accent2)

    # Date
    date_box = title_slide.shapes.add_textbox(
        PptxInches(1.5), PptxInches(4.8), PptxInches(5), PptxInches(0.5)
    )
    dp = date_box.text_frame.paragraphs[0]
    dp.text = datetime.now().strftime("%B %d, %Y")
    dp.font.size = PptxPt(14)
    dp.font.color.rgb = ACCENT2

    # ── Content slides ──
    slides_data = content.get("slides", [])
    total_slides = len(slides_data)

    for s_idx, sd in enumerate(slides_data):
        layout_name = sd.get("layout", "title_and_content")
        is_section = layout_name == "section"

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

        if is_section:
            # Section divider — dark background
            _pptx_set_bg(slide, *dark_bg)
            _add_pptx_shape_rect(slide, PptxInches(0), PptxInches(0),
                                 PptxInches(bar_w), slide_h, *accent)
            sect_box = slide.shapes.add_textbox(
                PptxInches(1.5), PptxInches(2.5), PptxInches(10), PptxInches(2.0)
            )
            sp = sect_box.text_frame.paragraphs[0]
            sp.text = sd.get("title", "")
            sp.font.size = PptxPt(min(title_font_sz, 36))
            sp.font.bold = True
            sp.font.color.rgb = WHITE
        else:
            # Normal content slide — light background
            _pptx_set_bg(slide, *light_bg)

            # Top accent bar
            _add_pptx_shape_rect(slide, PptxInches(0), PptxInches(0),
                                 slide_w, PptxInches(0.08), *accent)

            # Slide title area with dark background strip
            _add_pptx_shape_rect(slide, PptxInches(0), PptxInches(0.08),
                                 slide_w, PptxInches(1.1), *dark_bg)

            title_box = slide.shapes.add_textbox(
                PptxInches(0.8), PptxInches(0.2), PptxInches(11.5), PptxInches(0.9)
            )
            tp = title_box.text_frame.paragraphs[0]
            tp.text = sd.get("title", "")
            tp.font.size = PptxPt(min(title_font_sz - 12, 28))
            tp.font.bold = True
            tp.font.color.rgb = WHITE

            # Bullet content area
            bullets = sd.get("bullets", [])
            if bullets:
                body_box = slide.shapes.add_textbox(
                    PptxInches(1.0), PptxInches(1.6),
                    PptxInches(11.0), PptxInches(5.2)
                )
                btf = body_box.text_frame
                btf.word_wrap = True

                for i, bullet in enumerate(bullets):
                    if i == 0:
                        p = btf.paragraphs[0]
                    else:
                        p = btf.add_paragraph()

                    # Accent bullet marker
                    run_marker = p.add_run()
                    run_marker.text = f"{bullet_char}  "
                    run_marker.font.size = PptxPt(bullet_font_sz)
                    run_marker.font.color.rgb = ACCENT

                    run_text = p.add_run()
                    run_text.text = str(bullet)
                    run_text.font.size = PptxPt(body_font_sz)
                    run_text.font.color.rgb = DARK_TEXT

                    p.space_before = PptxPt(10)
                    p.space_after = PptxPt(4)

            # Slide number (bottom right)
            num_box = slide.shapes.add_textbox(
                PptxInches(12.0), PptxInches(7.0),
                PptxInches(1.0), PptxInches(0.4)
            )
            np_ = num_box.text_frame.paragraphs[0]
            np_.text = f"{s_idx + 1}/{total_slides}"
            np_.font.size = PptxPt(10)
            np_.font.color.rgb = RGBColor(120, 120, 120)
            from pptx.enum.text import PP_ALIGN
            np_.alignment = PP_ALIGN.RIGHT

        # Speaker notes
        notes = sd.get("notes", "")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(str(filepath))
    return str(filepath)


# ─── CSV ──────────────────────────────────────────────────────────────────

def generate_csv_file(content, filename=None):
    filepath = OUTPUT_DIR / (filename or f"data_{_ts()}.csv")
    with open(filepath, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        headers = content.get("headers", [])
        if headers:
            writer.writerow(headers)
        for row in content.get("rows", []):
            writer.writerow(row)
    return str(filepath)


# ─── JSON ─────────────────────────────────────────────────────────────────

def generate_json_file(content, filename=None):
    filepath = OUTPUT_DIR / (filename or f"data_{_ts()}.json")
    data = content.get("data", content)
    with open(filepath, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)
    return str(filepath)


# ─── Image (Chart) ───────────────────────────────────────────────────────

def generate_image(content, filename=None):
    filepath = OUTPUT_DIR / (filename or f"chart_{_ts()}.png")

    chart_type = content.get("chart_type", "bar")
    title = content.get("title", "Chart")
    xlabel = content.get("xlabel", "")
    ylabel = content.get("ylabel", "")
    data = content.get("data", {})
    labels = data.get("labels", [])
    datasets = data.get("datasets", [])

    fig, ax = plt.subplots(figsize=(10, 6))

    if chart_type == "pie" and datasets:
        values = datasets[0].get("values", [])
        ax.pie(values, labels=labels, autopct="%1.1f%%", startangle=90)
        ax.axis("equal")

    elif chart_type == "line":
        for ds in datasets:
            ax.plot(labels, ds.get("values", []), marker="o", label=ds.get("label", ""))
        ax.set_xlabel(xlabel)
        ax.set_ylabel(ylabel)
        if len(datasets) > 1:
            ax.legend()

    elif chart_type == "scatter":
        for ds in datasets:
            x_vals = ds.get("x", list(range(len(ds.get("values", [])))))
            ax.scatter(x_vals, ds.get("values", []), label=ds.get("label", ""))
        ax.set_xlabel(xlabel)
        ax.set_ylabel(ylabel)
        if len(datasets) > 1:
            ax.legend()

    else:  # bar (default)
        x = np.arange(len(labels))
        width = 0.8 / max(len(datasets), 1)
        for i, ds in enumerate(datasets):
            offset = (i - len(datasets) / 2 + 0.5) * width
            ax.bar(x + offset, ds.get("values", []), width, label=ds.get("label", ""))
        ax.set_xticks(x)
        ax.set_xticklabels(labels, rotation=45, ha="right")
        ax.set_xlabel(xlabel)
        ax.set_ylabel(ylabel)
        if len(datasets) > 1:
            ax.legend()

    ax.set_title(title)
    plt.tight_layout()
    plt.savefig(str(filepath), dpi=150, bbox_inches="tight")
    plt.close(fig)
    return str(filepath)


# ─── Dispatcher ───────────────────────────────────────────────────────────

GENERATORS = {
    "pdf": generate_pdf,
    "docx": generate_docx,
    "xlsx": generate_xlsx,
    "pptx": generate_pptx,
    "csv": generate_csv_file,
    "json": generate_json_file,
    "png": generate_image,
}


def generate_document(content_json, format_type, filename=None, design=None):
    """Generate a document in the specified format. Returns path to the file.

    Args:
        content_json: Structured content dict from LLM.
        format_type: Output format (pdf, docx, xlsx, pptx, csv, json, png).
        filename: Optional filename override.
        design: Optional design dict extracted from an uploaded file.
    """
    generator = GENERATORS.get(format_type.lower())
    if not generator:
        raise ValueError(f"Unsupported format: {format_type}. Supported: {list(GENERATORS.keys())}")
    # Pass design to generators that support it
    if format_type.lower() == "pptx" and design:
        return generator(content_json, filename, design=design)
    return generator(content_json, filename)
