"""Extract visual design properties from uploaded documents (PPTX, DOCX, PDF).

Returns a design dict that doc_generator can apply to new documents,
preserving the source file's color scheme, fonts, and layout style.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor


def _rgb_tuple(rgb_color):
    """Convert pptx RGBColor or similar to (r, g, b) tuple."""
    if rgb_color is None:
        return None
    if isinstance(rgb_color, RGBColor):
        return (rgb_color[0], rgb_color[1], rgb_color[2])
    if isinstance(rgb_color, (list, tuple)) and len(rgb_color) >= 3:
        return tuple(rgb_color[:3])
    return None


def _emu_to_inches(emu):
    if emu is None:
        return None
    return round(emu / 914400, 2)


def _safe_font_color(font):
    """Safely extract font color as (r, g, b) tuple."""
    try:
        if font.color and font.color.type is not None and font.color.rgb:
            return _rgb_tuple(font.color.rgb)
    except Exception:
        pass
    return None


def _safe_fill_color(fill):
    """Safely extract solid fill color as (r, g, b) tuple."""
    try:
        if fill.type is not None:
            fill_type_name = str(fill.type)
            if "SOLID" in fill_type_name:
                if fill.fore_color and fill.fore_color.type is not None:
                    return _rgb_tuple(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def extract_pptx_design(file_path):
    """Extract design properties from an existing PPTX file.

    Returns a dict with:
        slide_width, slide_height: in inches
        colors: dict of discovered colors (dark_bg, light_bg, accent, accent2, title_text, body_text)
        fonts: dict of font info (title_name, title_size, body_name, body_size, bullet_size)
        has_accent_bar: bool
        accent_bar_width: float (inches)
        bullet_char: str
    """
    prs = Presentation(str(file_path))

    design = {
        "source": Path(file_path).name,
        "slide_width": _emu_to_inches(prs.slide_width),
        "slide_height": _emu_to_inches(prs.slide_height),
        "colors": {},
        "fonts": {},
        "has_accent_bar": False,
        "accent_bar_width": 0.4,
        "bullet_char": "\u25A0",
    }

    # Collect all colors and fonts from slides
    bg_colors = []
    text_colors = []
    shape_fill_colors = []
    font_names = []
    title_font_sizes = []
    body_font_sizes = []

    for slide_idx, slide in enumerate(prs.slides):
        # Background color
        try:
            bg_fill = slide.background.fill
            bg_c = _safe_fill_color(bg_fill)
            if bg_c:
                bg_colors.append((slide_idx, bg_c))
        except Exception:
            pass

        for shape in slide.shapes:
            # Shape fill colors (rectangles, accent bars, etc.)
            try:
                if shape.has_text_frame:
                    pass  # check text below
                fill_c = _safe_fill_color(shape.fill)
                if fill_c:
                    # Detect accent bars (narrow rectangles)
                    w = _emu_to_inches(shape.width)
                    h = _emu_to_inches(shape.height)
                    if w and h:
                        if w < 1.0 and h > 2.0:
                            design["has_accent_bar"] = True
                            design["accent_bar_width"] = w
                        elif h < 0.15 and w > 3.0:
                            design["has_accent_bar"] = True
                    shape_fill_colors.append(fill_c)
            except Exception:
                pass

            # Text properties
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        text = run.text.strip()
                        if not text:
                            continue
                        f = run.font
                        fc = _safe_font_color(f)
                        if fc:
                            text_colors.append(fc)
                        if f.name:
                            font_names.append(f.name)
                        if f.size:
                            pt_size = f.size.pt if hasattr(f.size, "pt") else round(f.size / 12700)
                            if pt_size >= 24:
                                title_font_sizes.append(pt_size)
                                if f.name:
                                    design["fonts"]["title_name"] = f.name
                            elif pt_size >= 10:
                                body_font_sizes.append(pt_size)
                                if f.name:
                                    design["fonts"]["body_name"] = f.name

    # ── Analyze collected colors to identify design roles ──

    # Dark backgrounds (title slide, section dividers)
    dark_bgs = [c for _, c in bg_colors if _brightness(c) < 100]
    light_bgs = [c for _, c in bg_colors if _brightness(c) >= 180]

    if dark_bgs:
        design["colors"]["dark_bg"] = _most_common(dark_bgs)
    if light_bgs:
        design["colors"]["light_bg"] = _most_common(light_bgs)

    # Accent colors from shapes (not black/white/gray)
    accent_candidates = [c for c in shape_fill_colors if _is_chromatic(c)]
    if accent_candidates:
        sorted_by_freq = _sorted_by_frequency(accent_candidates)
        design["colors"]["accent"] = sorted_by_freq[0]
        if len(sorted_by_freq) > 1:
            design["colors"]["accent2"] = sorted_by_freq[1]

    # Title text color (usually white on dark bg or dark on light bg)
    white_texts = [c for c in text_colors if _brightness(c) > 220]
    dark_texts = [c for c in text_colors if _brightness(c) < 80]
    if white_texts:
        design["colors"]["title_text"] = (255, 255, 255)
    elif dark_texts:
        design["colors"]["title_text"] = _most_common(dark_texts)

    # Body text color
    mid_dark_texts = [c for c in text_colors if 20 < _brightness(c) < 120]
    if mid_dark_texts:
        design["colors"]["body_text"] = _most_common(mid_dark_texts)
    elif dark_texts:
        design["colors"]["body_text"] = _most_common(dark_texts)

    # Font sizes
    if title_font_sizes:
        design["fonts"]["title_size"] = round(sum(title_font_sizes) / len(title_font_sizes))
    if body_font_sizes:
        design["fonts"]["body_size"] = round(sum(body_font_sizes) / len(body_font_sizes))
        design["fonts"]["bullet_size"] = design["fonts"]["body_size"]

    return design


def _brightness(rgb_tuple):
    """Perceived brightness (0-255)."""
    r, g, b = rgb_tuple
    return 0.299 * r + 0.587 * g + 0.114 * b


def _is_chromatic(rgb_tuple):
    """Return True if the color is not black, white, or gray."""
    r, g, b = rgb_tuple
    spread = max(r, g, b) - min(r, g, b)
    return spread > 30  # has enough color saturation


def _most_common(color_list):
    """Return the most common color from a list, grouping similar ones."""
    if not color_list:
        return None
    # Simple frequency count
    from collections import Counter
    counts = Counter(color_list)
    return counts.most_common(1)[0][0]


def _sorted_by_frequency(color_list):
    """Return unique colors sorted by frequency (most common first)."""
    from collections import Counter
    counts = Counter(color_list)
    return [c for c, _ in counts.most_common()]


def extract_design(file_path):
    """Extract design from any supported file type. Returns design dict or None."""
    ext = Path(file_path).suffix.lower()
    if ext == ".pptx":
        return extract_pptx_design(file_path)
    # Future: .docx, .pdf design extraction
    return None


def describe_design(design):
    """Create a human-readable summary of the design for LLM context."""
    if not design:
        return ""

    parts = [f"Design extracted from: {design.get('source', 'uploaded file')}"]

    colors = design.get("colors", {})
    if colors.get("dark_bg"):
        parts.append(f"  Dark background: RGB{colors['dark_bg']}")
    if colors.get("light_bg"):
        parts.append(f"  Light background: RGB{colors['light_bg']}")
    if colors.get("accent"):
        parts.append(f"  Primary accent: RGB{colors['accent']}")
    if colors.get("accent2"):
        parts.append(f"  Secondary accent: RGB{colors['accent2']}")

    fonts = design.get("fonts", {})
    if fonts.get("title_name"):
        parts.append(f"  Title font: {fonts['title_name']} {fonts.get('title_size', '')}pt")
    if fonts.get("body_name"):
        parts.append(f"  Body font: {fonts['body_name']} {fonts.get('body_size', '')}pt")

    if design.get("has_accent_bar"):
        parts.append("  Has accent/decorative bars")

    return "\n".join(parts)
